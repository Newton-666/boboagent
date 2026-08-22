"""读取本地文件内容（支持文件和目录）"""

import os
from pathlib import Path
from core.file_safety import safe_read_check

TOOL_NAME = "read_local_file"

# 读取目录时每个文件最多显示的行数
DIR_PREVIEW_LINES = 30

# TICKET-VISION-INPUT：图片扩展名 → 走视觉描述分支（真正让视觉模型看图，非元信息）
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp'}


def _describe_image(filepath: str, max_chars: int = 40000) -> str:
    """图片 → base64 data URL → 视觉描述（真正看图，非读取元信息）。

    TICKET-VISION-INPUT：支持 vision 的模型 → 构造多模态消息（image_url + base64）
    调视觉模型返回描述；非 vision 模型 → 明确报错（不静默）。
    依赖：core.provider.supports_vision 判定 provider/model 是否支持图像输入。
    """
    import base64
    import mimetypes
    from core.provider import resolve_provider, supports_vision

    try:
        with open(filepath, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("ascii")
    except OSError as e:
        return f"错误: 读取图片失败: {e}"

    mime = mimetypes.guess_type(filepath)[0] or "image/png"
    data_url = f"data:{mime};base64,{b64}"

    cfg = resolve_provider()
    model = cfg.get("model", "")
    if not supports_vision(cfg.get("name", ""), model):
        return (f"错误: 当前模型 {model}（provider={cfg.get('name')}）不支持图像输入 vision。"
                f"请切换到支持 vision 的模型（如 deepseek-v4-flash-vision-exp）。")

    from core.llm_caller import create_llm_caller
    llm = create_llm_caller(cfg.get("api_key", ""), cfg.get("base_url", ""), model)
    messages = [{
        "role": "user",
        "content": [
            {"type": "text", "text": "请描述这张图片的内容（颜色、物体、文字、场景）。用中文简洁回答。"},
            {"type": "image_url", "image_url": {"url": data_url}},
        ],
    }]
    try:
        resp = llm(messages, use_tools=False, max_tokens=1024, thinking_disabled=True)
    except Exception as e:
        return f"错误: 视觉模型调用失败: {e}"
    if isinstance(resp, dict) and resp.get("error"):
        return f"错误: 视觉模型返回错误: {resp['error']}"
    content = ""
    try:
        content = resp["choices"][0]["message"].get("content", "") or ""
    except (KeyError, IndexError, TypeError):
        content = ""
    if not content:
        return f"{filepath}\n\n[视觉描述] （模型未返回内容）"
    return f"{filepath}\n\n[视觉描述] {content[:max_chars]}" + \
        (f"\n... (内容已截断，共 {len(content)} 字符)" if len(content) > max_chars else "")


def _read_single_file(filepath: str, max_chars: int = 40000,
                      offset: int = 0, limit: int | None = None) -> str:
    """读取单个文件内容

    Args:
        filepath: 文件路径
        max_chars: 最大返回字符数（默认 40000）
        offset: 从第几行开始读（0 表示开头）
        limit: 最多读几行（None 表示全部）
    """
    path = Path(filepath).expanduser()

    if not path.exists():
        return f"错误: 文件不存在: {filepath}"

    ext = path.suffix.lower()

    try:
        if ext == '.pdf':
            try:
                import pypdf
                reader = pypdf.PdfReader(str(path))
                text = "\n".join(page.extract_text() or "" for page in reader.pages)
                return f"{filepath}\n\n{text[:max_chars]}" + (f"\n... (共 {len(text)} 字符)" if len(text) > max_chars else "")
            except ImportError:
                return "错误: 请安装 pypdf: pip install pypdf"
        elif ext in ['.docx', '.doc']:
            try:
                import docx
                doc = docx.Document(str(path))
                text = "\n".join(p.text for p in doc.paragraphs)
                return f"{filepath}\n\n{text[:max_chars]}" + (f"\n... (共 {len(text)} 字符)" if len(text) > max_chars else "")
            except ImportError:
                return "错误: 请安装 python-docx: pip install python-docx"
        elif ext in ['.pptx', '.ppt']:
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                text = "\n".join(text_parts)
                return f"{filepath}\n\n{text[:max_chars]}" + (f"\n... (共 {len(text)} 字符)" if len(text) > max_chars else "")
            except ImportError:
                return "错误: 请安装 python-pptx: pip install python-pptx"
        elif ext in ['.xlsx', '.xls']:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
                text_parts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    rows = []
                    for row in ws.iter_rows(values_only=True):
                        row_text = " | ".join(str(c) for c in row if c is not None)
                        if row_text.strip():
                            rows.append(row_text)
                    if rows:
                        text_parts.append(f"[{sheet_name}]\n" + "\n".join(rows[:50]))
                        if len(rows) > 50:
                            text_parts.append(f"... (共 {len(rows)} 行)")
                wb.close()
                text = "\n\n".join(text_parts)
                return f"{filepath}\n\n{text[:max_chars]}" + (f"\n... (共 {len(text)} 字符)" if len(text) > max_chars else "")
            except ImportError:
                return "错误: 请安装 openpyxl: pip install openpyxl"
        elif ext in ['.md', '.txt', '.py', '.json', '.yaml', '.yml', '.html', '.css', '.js', '.sh']:
            content = path.read_text(encoding='utf-8')
        else:
            content = path.read_text(encoding='utf-8', errors='ignore')

        # ── 分页：offset + limit ──
        if offset > 0 or limit is not None:
            lines = content.split('\n')
            total_lines = len(lines)
            start = offset
            end = (start + limit) if limit is not None else total_lines
            selected = lines[start:end]
            content = '\n'.join(selected)
            if offset > 0 or limit is not None:
                header = f"[行 {start+1}-{min(end, total_lines)} / 共 {total_lines} 行]\n"
                content = header + content

        if len(content) > max_chars:
            content = content[:max_chars] + f"\n... (内容已截断，共 {len(content)} 字符)"

        return f"{filepath}\n\n{content}"
    except Exception as e:
        return f"错误: 读取失败: {str(e)}"


def _read_directory(dirpath: str) -> str:
    """读取目录结构，返回每个文件的摘要"""
    path = Path(dirpath).expanduser()

    if not path.exists():
        return f"错误: 目录不存在: {dirpath}"
    if not path.is_dir():
        return _read_single_file(dirpath)

    result = []
    result.append(f"目录: {dirpath}")
    result.append("")

    # 收集所有文件
    files = []
    for f in sorted(path.iterdir()):
        if f.name.startswith('.'):
            continue
        if f.is_file():
            size = f.stat().st_size
            files.append((f.name, size, f))

    result.append(f"共 {len(files)} 个文件")
    result.append("")

    for name, size, fpath in files:
        size_str = f"{size}B" if size < 1024 else f"{size/1024:.1f}KB"
        result.append(f"  {name} ({size_str})")

        # 读取前几行作为预览
        try:
            ext = fpath.suffix.lower()
            if ext in ['.md', '.txt', '.py', '.json', '.yaml', '.yml', '.html', '.css', '.js', '.sh']:
                lines = fpath.read_text(encoding='utf-8', errors='ignore').split('\n')
                preview_lines = lines[:DIR_PREVIEW_LINES]
                for line in preview_lines:
                    if line.strip():
                        result.append(f"    {line[:100]}")
                if len(lines) > DIR_PREVIEW_LINES:
                    result.append(f"    ... (共 {len(lines)} 行)")
        except Exception:
            pass
        result.append("")

    return '\n'.join(result)


def execute(filepath: str, max_chars: int = 40000,
            offset: int = 0, limit: int = None) -> str:
    """读取本地文件或目录内容

    Args:
        filepath: 文件或目录路径
        max_chars: 最大返回字符数（默认 40000）
        offset: 从第几行开始读（默认 0 = 开头）。用于大文件分页读取
        limit: 最多读取的行数（默认 None = 全部）。配合 offset 实现分页
    """
    path = Path(filepath).expanduser()

    if not path.exists():
        return f"错误: 路径不存在: {filepath}"

    # 安全: 二进制文件检测（图片走视觉分支，绕过二进制拦截 TICKET-VISION-INPUT）
    if path.is_file():
        if path.suffix.lower() in IMAGE_EXTENSIONS:
            return _describe_image(str(path), max_chars)
        warning = safe_read_check(str(path))
        if warning:
            return warning

    if path.is_dir():
        return _read_directory(filepath)
    else:
        return _read_single_file(filepath, max_chars, offset, limit)


TOOL_FUNC = execute
TOOL_SCHEMA = {
    "type": "function",
    "function": {
        "name": TOOL_NAME,
        "description": "读取本地文件或目录内容（默认上限 40000 字符）。支持 .md, .txt, .py, .pdf, .docx 等。大文件可用 offset+limit 分页读取，防止撑爆上下文。目录返回结构预览。",
        "parameters": {"type": "object", "properties": {
            "filepath": {"type": "string", "description": "要读取的文件绝对路径"},
            "max_chars": {"type": "integer", "description": "最大返回字符数，超出的内容截断"},
            "offset": {"type": "integer", "description": "从第几行开始读（0=开头），大文件分页用"},
            "limit": {"type": "integer", "description": "最多读取行数，配合 offset 分页"}
        }, "required": ["filepath"]}
    }
}
def register(reg): reg(TOOL_NAME, TOOL_FUNC, TOOL_SCHEMA)
